"""PowerPoint report generation with chart images and formatted tables."""

import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ils_kickoff.formatting import format_value, is_grand_total_row

logger = logging.getLogger(__name__)

# Style constants
NAVY = RGBColor(0x1B, 0x36, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_FILL = RGBColor(0xCC, 0xE5, 0xFF)
TOTAL_FILL = RGBColor(0xF0, 0xF0, 0xF0)
ZEBRA_FILL = RGBColor(0xFA, 0xFA, 0xFA)


def _add_title_slide(prs: Presentation, settings) -> None:
    """Add a navy title slide with client name and date."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = "ILS Kickoff Analysis"
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = NAVY

    # Subtitle
    subtitle = slide.placeholders[1]
    now = datetime.now()
    subtitle.text = f"{settings.client_name} | {now.strftime('%B %Y')}"
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _format_table(table, df: pd.DataFrame) -> None:
    """Apply formatting to a PowerPoint table."""
    rows, cols = df.shape

    # Header row
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r_idx, (_, row) in enumerate(df.iterrows()):
        is_total = is_grand_total_row(row)
        is_odd = r_idx % 2 == 1

        for c_idx, col_name in enumerate(df.columns):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = format_value(row[col_name], col_name)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

            if is_total:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = TOTAL_FILL
            elif is_odd:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ZEBRA_FILL


def _add_analysis_slide(
    prs: Presentation,
    title: str,
    df: pd.DataFrame,
    chart_png: bytes | None = None,
) -> None:
    """Add a slide with optional chart image and data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.top = Inches(0.2)
    title_shape.left = Inches(0.3)
    title_shape.width = Inches(9.2)
    title_shape.height = Inches(0.7)
    tf = title_shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = NAVY

    # Chart image (if available)
    table_top = Inches(1.0)
    if chart_png:
        chart_stream = BytesIO(chart_png)
        slide.shapes.add_picture(
            chart_stream,
            left=Inches(0.5),
            top=Inches(1.0),
            width=Inches(5.5),
        )
        table_top = Inches(4.3)

    # Data table
    n_rows = min(df.shape[0] + 1, 26)  # Cap at 25 data rows + header
    n_cols = df.shape[1]

    # Adaptive sizing
    font_size = Pt(8) if df.shape[0] <= 15 else Pt(7)
    row_height = 0.25
    table_height = min(row_height * n_rows + 0.2, 5.0)
    col_width = min(9.0 / n_cols, 1.5)
    table_width = min(col_width * n_cols + 0.5, 9.5)

    # If too many rows and we have a chart, use separate slide for table
    display_df = df.head(25) if df.shape[0] > 25 and chart_png else df

    table_shape = slide.shapes.add_table(
        display_df.shape[0] + 1, n_cols,
        Inches(0.3), table_top,
        Inches(table_width), Inches(table_height),
    )
    _format_table(table_shape.table, display_df)

    # If we truncated, add overflow slide with full table
    if df.shape[0] > 25 and chart_png:
        overflow_slide = prs.slides.add_slide(prs.slide_layouts[5])
        overflow_title = overflow_slide.shapes.title
        overflow_title.text = f"{title} (Full Table)"
        overflow_title.top = Inches(0.2)
        overflow_title.left = Inches(0.3)
        overflow_title.width = Inches(9.2)
        overflow_title.height = Inches(0.5)
        for para in overflow_title.text_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.bold = True

        full_rows = df.shape[0] + 1
        full_table = overflow_slide.shapes.add_table(
            full_rows, n_cols,
            Inches(0.3), Inches(0.8),
            Inches(table_width), Inches(min(row_height * full_rows + 0.2, 6.5)),
        )
        _format_table(full_table.table, df)


def write_pptx_report(result, output_path: Path) -> None:
    """Write the complete PowerPoint report."""
    prs = Presentation()

    _add_title_slide(prs, result.settings)

    # Render chart PNGs
    chart_pngs: dict[str, bytes] = {}
    if result.charts:
        try:
            from ils_kickoff.charts import render_chart_png
            for name, fig in result.charts.items():
                try:
                    chart_pngs[name] = render_chart_png(fig, result.settings.charts)
                except Exception as e:
                    logger.warning("Chart PNG for '%s' failed: %s", name, e)
        except ImportError:
            pass

    # Add analysis slides
    for analysis in result.analyses:
        if analysis.error is not None or analysis.df.empty:
            continue

        _add_analysis_slide(
            prs,
            title=analysis.title,
            df=analysis.df,
            chart_png=chart_pngs.get(analysis.name),
        )

    prs.save(output_path)
    logger.info("PowerPoint saved: %s", output_path)
