#!/usr/bin/env python3
"""
ILS Kickoff Presentation Generator (Cleaned)
=============================================
Loads client OD/NSF data, runs analyses, and generates a PowerPoint deck + Excel report.
"""

# ═══════════════════════════════════════════════════════════════════
# 1. IMPORTS & CONFIGURATION
# ═══════════════════════════════════════════════════════════════════
import pandas as pd
import numpy as np
from pathlib import Path
from datetime import datetime
import calendar
import re
import os
import warnings

import matplotlib
matplotlib.use("Agg")  # Non-interactive backend
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from openpyxl import Workbook, load_workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows

warnings.filterwarnings("ignore")
pd.set_option("display.max_columns", 100)
pd.set_option("display.width", 120)
pd.options.mode.chained_assignment = None

plt.rcParams["figure.dpi"] = 120
plt.rcParams["savefig.bbox"] = "tight"
plt.rcParams["savefig.dpi"] = 300
plt.rcParams["figure.figsize"] = (10, 6)

# ═══════════════════════════════════════════════════════════════════
# 2. DATA LOADING & CLEANING
# ═══════════════════════════════════════════════════════════════════
DATA_PATH = Path("1774_INB_OD_Tran_Combo_20260203__1_.csv")

print(f"Loading: {DATA_PATH.name}")
df = pd.read_csv(DATA_PATH, encoding="utf-8-sig", low_memory=False)
print(f"Shape: {df.shape[0]:,} rows × {df.shape[1]} columns")

# ── Rename columns to canonical names ──
RENAME_MAP = {
    "TOTALITEMS": "Total Items",
    "PaidItems": "Paid Items",
    "ReturnedItems": "Returned Items",
    "ODLimit": "OD Limit",
    "ODStatus": "OD Status",
    "ProdCode": "Product Code",
    "BusinessFlag": "Business Flag",
    "AccountStatus": "Account Status",
    "RegEValue": "Reg E Flag",
    "OpenDate": "Open Date",
    "AvgColBal": "Avg Bal",
    "DepositAmount": "Deposit Amount",
    "DepositCount": "Deposit Count",
    "swipes": "Swipes",
    "spend": "Spend",
}
df.rename(columns=RENAME_MAP, inplace=True)

# ── Fix DepositCount: extract first value before embedded tabs ──
df["Deposit Count"] = (
    df["Deposit Count"]
    .astype(str)
    .str.split("\t")
    .str[0]
    .pipe(pd.to_numeric, errors="coerce")
    .fillna(0)
    .astype(int)
)

# ── Fix Returned Items (all NaN → 0) ──
df["Returned Items"] = df["Returned Items"].fillna(0).astype(int)

# ── Parse Open Date ──
df["Open Date"] = pd.to_datetime(df["Open Date"], errors="coerce")
df["Year Opened"] = df["Open Date"].dt.year

# ── Ensure numerics ──
for col in ["Total Items", "Paid Items", "OD Limit", "Avg Bal", "Deposit Amount", "Swipes", "Spend"]:
    df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

print(f"Columns: {list(df.columns)}")
print(f"Account Status values: {sorted(df['Account Status'].unique())}")
print(f"Business Flag values: {sorted(df['Business Flag'].unique())}")
print(f"Reg E Flag values: {sorted(df['Reg E Flag'].dropna().unique())}")

# ═══════════════════════════════════════════════════════════════════
# 3. HELPER / FORMATTING FUNCTIONS
# ═══════════════════════════════════════════════════════════════════

def add_grand_total(summary_df, label_col, label="Grand Total"):
    """Add a Grand Total row to a summary DataFrame."""
    totals = {}
    for col in summary_df.columns:
        if col == label_col:
            totals[col] = label
        elif "%" in col or "Ratio" in col.lower() or "Avg" in col or "Average" in col or "Med" in col:
            # Recalculate percentage/ratio from totals
            totals[col] = np.nan  # placeholder
        else:
            totals[col] = summary_df[col].sum()
    return totals


def format_ppt_table(table, data, header_fill=RGBColor(204, 229, 255)):
    """Apply formatting to a PowerPoint table."""
    rows, cols = data.shape

    # Header row
    for col_idx, col_name in enumerate(data.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r_idx, (_, row) in enumerate(data.iterrows()):
        for c_idx, col_name in enumerate(data.columns):
            cell = table.cell(r_idx + 1, c_idx)
            val = row[col_name]

            # Format value
            if pd.isna(val) or val == "":
                cell.text = ""
            elif isinstance(val, (int, np.integer)):
                cell.text = f"{val:,}"
            elif isinstance(val, (float, np.floating)):
                if "%" in col_name:
                    cell.text = f"{val:.1f}%"
                elif "Ratio" in col_name:
                    cell.text = f"{val:.2f}"
                elif "$$" in col_name or "Limit" in col_name or "Dep" in col_name.split("/")[0]:
                    cell.text = f"${val:,.0f}" if abs(val) >= 100 else f"${val:.2f}"
                elif "Avg" in col_name or "Average" in col_name or "Med" in col_name:
                    cell.text = f"{val:.1f}"
                else:
                    cell.text = f"{val:,.0f}" if abs(val) >= 10 else f"{val:.2f}"
            else:
                cell.text = str(val)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.alignment = PP_ALIGN.CENTER

            # Bold Grand Total rows
            first_val = str(row.iloc[0]).lower()
            if "total" in first_val or "grand" in first_val:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)


def add_slide_with_table(prs, title, data_df):
    """Add a slide with a formatted table to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Set title with constrained size
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.top = Inches(0.2)
    title_shape.left = Inches(0.3)
    title_shape.width = Inches(9.2)
    title_shape.height = Inches(0.8)
    tf = title_shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.font.size = Pt(24)
        para.font.bold = True

    n_rows = data_df.shape[0] + 1  # +1 for header
    n_cols = data_df.shape[1]

    # Calculate dynamic sizing
    row_height = 0.28
    table_height = min(row_height * n_rows + 0.3, 5.5)
    col_width = min(9.0 / n_cols, 1.5)
    table_width = min(col_width * n_cols + 0.5, 9.5)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.3), Inches(1.2),
        Inches(table_width), Inches(table_height)
    )
    format_ppt_table(table_shape.table, data_df)
    return slide


def save_chart(fig, filename):
    """Save a matplotlib figure and close it."""
    fig.savefig(filename, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return filename


# ═══════════════════════════════════════════════════════════════════
# 4. RUN ALL ANALYSES
# ═══════════════════════════════════════════════════════════════════

# Common filters
OPEN_MASK = df["Account Status"] == "O"
PERSONAL_OPEN = OPEN_MASK & (df["Business Flag"] == "P")
BUSINESS_OPEN = OPEN_MASK & (df["Business Flag"] == "B")

NSF_BINS = [-1, 0, 6, 12, 24, 36, 48, float("inf")]
NSF_LABELS = ["0", "1–6", "7–12", "13–24", "25–36", "37–48", "49+"]

DEP_BINS = [-1, 0, 1, 2, 5, 10, float("inf")]
DEP_LABELS = ["0", "1", "2", "3–5", "6–10", "10+"]

chart_files = []  # Track temp chart files for cleanup

# ── ANALYSIS 1: Account Status Summary (All Accounts) ──
print("\n── Analysis 1: Account Status Summary ──")
stat_code_summary = df.groupby("Account Status").agg(
    **{
        "# of Accounts": ("AcctNo", "count"),
        "# of Items": ("Total Items", "sum"),
        "# of Items Paid": ("Paid Items", "sum"),
    }
).reset_index()

stat_code_summary["% of Accounts"] = (stat_code_summary["# of Accounts"] / stat_code_summary["# of Accounts"].sum() * 100).round(2)
stat_code_summary["% of Items"] = (stat_code_summary["# of Items"] / stat_code_summary["# of Items"].sum() * 100).round(2)
stat_code_summary["Pay Ratio"] = np.where(
    stat_code_summary["# of Items"] > 0,
    (stat_code_summary["# of Items Paid"] / stat_code_summary["# of Items"]).round(2), 0
)

# Grand Total
ti = stat_code_summary["# of Items"].sum()
tp = stat_code_summary["# of Items Paid"].sum()
gt = pd.DataFrame([{
    "Account Status": "Grand Total",
    "# of Accounts": stat_code_summary["# of Accounts"].sum(),
    "# of Items": ti, "# of Items Paid": tp,
    "% of Accounts": 100.0, "% of Items": 100.0,
    "Pay Ratio": round(tp / ti, 2) if ti > 0 else 0,
}])
stat_code_summary = pd.concat([stat_code_summary, gt], ignore_index=True)
stat_code_summary = stat_code_summary[["Account Status", "# of Accounts", "% of Accounts", "# of Items", "% of Items", "# of Items Paid", "Pay Ratio"]]
print(f"  {stat_code_summary.shape[0]} rows")

# ── ANALYSIS 2: Account Type (Open Accounts) ──
print("── Analysis 2: Account Type (Open Accounts) ──")
df_open = df[OPEN_MASK].copy()

acct_type_summary = df_open.groupby("Business Flag").agg(
    **{
        "# of Accounts": ("AcctNo", "count"),
        "# of Items": ("Total Items", "sum"),
        "# of Paid Items": ("Paid Items", "sum"),
    }
).reset_index()

acct_type_summary["% of Accounts"] = (acct_type_summary["# of Accounts"] / acct_type_summary["# of Accounts"].sum() * 100).round(2)
acct_type_summary["% of Items"] = (acct_type_summary["# of Items"] / acct_type_summary["# of Items"].sum() * 100).round(2)
acct_type_summary["Pay Ratio"] = np.where(
    acct_type_summary["# of Items"] > 0,
    (acct_type_summary["# of Paid Items"] / acct_type_summary["# of Items"]).round(2), 0
)
acct_type_summary["Business Flag"] = acct_type_summary["Business Flag"].map({"P": "Personal", "B": "Business"}).fillna(acct_type_summary["Business Flag"])

# Grand Total
ti2 = acct_type_summary["# of Items"].sum()
tp2 = acct_type_summary["# of Paid Items"].sum()
gt2 = pd.DataFrame([{
    "Business Flag": "Grand Total",
    "# of Accounts": acct_type_summary["# of Accounts"].sum(),
    "# of Items": ti2, "# of Paid Items": tp2,
    "% of Accounts": 100.0, "% of Items": 100.0,
    "Pay Ratio": round(tp2 / ti2, 2) if ti2 > 0 else 0,
}])
acct_type_summary = pd.concat([acct_type_summary, gt2], ignore_index=True)
print(f"  {acct_type_summary.shape[0]} rows")


# ── ANALYSIS 3: Personal Deposit Distribution ──
print("── Analysis 3: Personal Deposit Distribution ──")
personal_df = df[PERSONAL_OPEN].copy()
personal_df["Deposit Bin"] = pd.cut(personal_df["Deposit Count"], bins=DEP_BINS, labels=DEP_LABELS)

personal_deposit_summary = personal_df.groupby("Deposit Bin", observed=True).agg(
    Accounts=("AcctNo", "count"),
    **{"Avg $$ Deposits": ("Deposit Amount", "mean")},
    **{"Avg # Deposits": ("Deposit Count", "mean")},
).reset_index()
personal_deposit_summary["% of Accounts"] = (personal_deposit_summary["Accounts"] / personal_deposit_summary["Accounts"].sum() * 100).round(2)
personal_deposit_summary.rename(columns={"Deposit Bin": "Deposit Bin"}, inplace=True)

# Grand Total
gt3 = pd.DataFrame([{
    "Deposit Bin": "Grand Total",
    "Accounts": personal_deposit_summary["Accounts"].sum(),
    "Avg $$ Deposits": personal_df["Deposit Amount"].mean(),
    "Avg # Deposits": personal_df["Deposit Count"].mean(),
    "% of Accounts": 100.0,
}])
personal_deposit_summary = pd.concat([personal_deposit_summary, gt3], ignore_index=True)
print(f"  {personal_deposit_summary.shape[0]} rows")


# ── ANALYSIS 4: Business Deposit Distribution ──
print("── Analysis 4: Business Deposit Distribution ──")
business_df = df[BUSINESS_OPEN].copy()
business_df["Deposit Bin"] = pd.cut(business_df["Deposit Count"], bins=DEP_BINS, labels=DEP_LABELS)

business_deposit_summary = business_df.groupby("Deposit Bin", observed=True).agg(
    Accounts=("AcctNo", "count"),
    **{"Avg $$ Deposits": ("Deposit Amount", "mean")},
    **{"Avg # Deposits": ("Deposit Count", "mean")},
).reset_index()
business_deposit_summary["% of Accounts"] = (business_deposit_summary["Accounts"] / business_deposit_summary["Accounts"].sum() * 100).round(2)

gt4 = pd.DataFrame([{
    "Deposit Bin": "Grand Total",
    "Accounts": business_deposit_summary["Accounts"].sum(),
    "Avg $$ Deposits": business_df["Deposit Amount"].mean(),
    "Avg # Deposits": business_df["Deposit Count"].mean(),
    "% of Accounts": 100.0,
}])
business_deposit_summary = pd.concat([business_deposit_summary, gt4], ignore_index=True)
print(f"  {business_deposit_summary.shape[0]} rows")


# ── ANALYSIS 5: Personal NSF Stratification (Volume) ──
print("── Analysis 5: Personal NSF Stratification ──")
personal_nsf = df[PERSONAL_OPEN].copy()
personal_nsf["NSF Bin"] = pd.cut(personal_nsf["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_strat1 = personal_nsf.groupby("NSF Bin", observed=True).agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum")}
).reset_index()
nsf_strat1["NSF Bin"] = nsf_strat1["NSF Bin"].astype(str)

ta5 = nsf_strat1["# of Accounts"].sum()
ti5 = nsf_strat1["Total OD/NSF Items"].sum()
nsf_strat1["% of Accounts"] = (nsf_strat1["# of Accounts"] / ta5 * 100).round(2)
nsf_strat1["% of Items Presented"] = (nsf_strat1["Total OD/NSF Items"] / ti5 * 100).round(2) if ti5 > 0 else 0

gt5 = pd.DataFrame([{"NSF Bin": "Grand Total", "# of Accounts": ta5, "% of Accounts": 100.0, "Total OD/NSF Items": ti5, "% of Items Presented": 100.0}])
nsf_strat1 = pd.concat([nsf_strat1, gt5], ignore_index=True)
nsf_strat1 = nsf_strat1[["NSF Bin", "# of Accounts", "% of Accounts", "Total OD/NSF Items", "% of Items Presented"]]
print(f"  {nsf_strat1.shape[0]} rows")


# ── ANALYSIS 6: Business NSF Stratification (Volume) ──
print("── Analysis 6: Business NSF Stratification ──")
business_nsf = df[BUSINESS_OPEN].copy()
business_nsf["NSF Bin"] = pd.cut(business_nsf["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_strat_biz = business_nsf.groupby("NSF Bin", observed=True).agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum")}
).reset_index()
nsf_strat_biz["NSF Bin"] = nsf_strat_biz["NSF Bin"].astype(str)

ta6 = nsf_strat_biz["# of Accounts"].sum()
ti6 = nsf_strat_biz["Total OD/NSF Items"].sum()
nsf_strat_biz["% of Accounts"] = (nsf_strat_biz["# of Accounts"] / ta6 * 100).round(2)
nsf_strat_biz["% of Items Presented"] = (nsf_strat_biz["Total OD/NSF Items"] / ti6 * 100).round(2) if ti6 > 0 else 0

gt6 = pd.DataFrame([{"NSF Bin": "Grand Total", "# of Accounts": ta6, "% of Accounts": 100.0, "Total OD/NSF Items": ti6, "% of Items Presented": 100.0}])
nsf_strat_biz = pd.concat([nsf_strat_biz, gt6], ignore_index=True)
nsf_strat_biz = nsf_strat_biz[["NSF Bin", "# of Accounts", "% of Accounts", "Total OD/NSF Items", "% of Items Presented"]]
print(f"  {nsf_strat_biz.shape[0]} rows")


# ── ANALYSIS 7: Personal NSF + Pay Ratio ──
print("── Analysis 7: Personal NSF + Pay Ratio ──")
nsf_strat_personal_pay = df[PERSONAL_OPEN].copy()
nsf_strat_personal_pay["NSF Bin"] = pd.cut(nsf_strat_personal_pay["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_pay_p = nsf_strat_personal_pay.groupby("NSF Bin", observed=True).agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum"), "# of Items Paid": ("Paid Items", "sum")}
).reset_index()
nsf_pay_p["NSF Bin"] = nsf_pay_p["NSF Bin"].astype(str)

ta7 = nsf_pay_p["# of Accounts"].sum()
ti7 = nsf_pay_p["Total OD/NSF Items"].sum()
tp7 = nsf_pay_p["# of Items Paid"].sum()
nsf_pay_p["% of Accounts"] = (nsf_pay_p["# of Accounts"] / ta7 * 100).round(2)
nsf_pay_p["% of Items Presented"] = (nsf_pay_p["Total OD/NSF Items"] / ti7 * 100).round(2) if ti7 > 0 else 0
nsf_pay_p["% Pay Rate"] = np.where(nsf_pay_p["Total OD/NSF Items"] > 0, (nsf_pay_p["# of Items Paid"] / nsf_pay_p["Total OD/NSF Items"] * 100).round(1), 0)

gt7 = pd.DataFrame([{
    "NSF Bin": "Grand Total", "# of Accounts": ta7, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti7, "% of Items Presented": 100.0,
    "# of Items Paid": tp7, "% Pay Rate": round(tp7 / ti7 * 100, 1) if ti7 > 0 else 0,
}])
nsf_pay_p = pd.concat([nsf_pay_p, gt7], ignore_index=True)
print(f"  {nsf_pay_p.shape[0]} rows")


# ── ANALYSIS 8: Business NSF + Pay Ratio ──
print("── Analysis 8: Business NSF + Pay Ratio ──")
nsf_strat_business_pay = df[BUSINESS_OPEN].copy()
nsf_strat_business_pay["NSF Bin"] = pd.cut(nsf_strat_business_pay["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_pay_b = nsf_strat_business_pay.groupby("NSF Bin", observed=True).agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum"), "# of Items Paid": ("Paid Items", "sum")}
).reset_index()
nsf_pay_b["NSF Bin"] = nsf_pay_b["NSF Bin"].astype(str)

ta8 = nsf_pay_b["# of Accounts"].sum()
ti8 = nsf_pay_b["Total OD/NSF Items"].sum()
tp8 = nsf_pay_b["# of Items Paid"].sum()
nsf_pay_b["% of Accounts"] = (nsf_pay_b["# of Accounts"] / ta8 * 100).round(2)
nsf_pay_b["% of Items Presented"] = (nsf_pay_b["Total OD/NSF Items"] / ti8 * 100).round(2) if ti8 > 0 else 0
nsf_pay_b["% Pay Rate"] = np.where(nsf_pay_b["Total OD/NSF Items"] > 0, (nsf_pay_b["# of Items Paid"] / nsf_pay_b["Total OD/NSF Items"] * 100).round(1), 0)

gt8 = pd.DataFrame([{
    "NSF Bin": "Grand Total", "# of Accounts": ta8, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti8, "% of Items Presented": 100.0,
    "# of Items Paid": tp8, "% Pay Rate": round(tp8 / ti8 * 100, 1) if ti8 > 0 else 0,
}])
nsf_pay_b = pd.concat([nsf_pay_b, gt8], ignore_index=True)
print(f"  {nsf_pay_b.shape[0]} rows")


# ── ANALYSIS 9: Personal NSF + Deposits + Swipes ──
print("── Analysis 9: Personal NSF + Deposits + Swipes ──")
p9 = df[PERSONAL_OPEN].copy()
p9["NSF Bin"] = pd.cut(p9["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_deps_p = p9.groupby("NSF Bin", observed=True).agg(
    **{
        "# of Accounts": ("AcctNo", "count"),
        "Total OD/NSF Items": ("Total Items", "sum"),
        "# of Items Paid": ("Paid Items", "sum"),
        "Avg # Dep/Month": ("Deposit Count", "mean"),
        "Avg $$ Dep/Month": ("Deposit Amount", "mean"),
        "Average of OD Limit": ("OD Limit", "mean"),
        "Average of Swipes": ("Swipes", "mean"),
    }
).reset_index()
nsf_deps_p["NSF Bin"] = nsf_deps_p["NSF Bin"].astype(str)

ta9 = nsf_deps_p["# of Accounts"].sum()
ti9 = nsf_deps_p["Total OD/NSF Items"].sum()
tp9 = nsf_deps_p["# of Items Paid"].sum()
nsf_deps_p["% of Accounts"] = (nsf_deps_p["# of Accounts"] / ta9 * 100).round(2)
nsf_deps_p["% of Items Presented"] = (nsf_deps_p["Total OD/NSF Items"] / ti9 * 100).round(2) if ti9 > 0 else 0
nsf_deps_p["% Pay Rate"] = np.where(nsf_deps_p["Total OD/NSF Items"] > 0, (nsf_deps_p["# of Items Paid"] / nsf_deps_p["Total OD/NSF Items"] * 100).round(1), 0)

gt9 = pd.DataFrame([{
    "NSF Bin": "Grand Total", "# of Accounts": ta9, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti9, "% of Items Presented": 100.0,
    "# of Items Paid": tp9, "% Pay Rate": round(tp9 / ti9 * 100, 1) if ti9 > 0 else 0,
    "Avg # Dep/Month": p9["Deposit Count"].mean(),
    "Avg $$ Dep/Month": p9["Deposit Amount"].mean(),
    "Average of OD Limit": p9["OD Limit"].mean(),
    "Average of Swipes": p9["Swipes"].mean(),
}])
nsf_deps_p = pd.concat([nsf_deps_p, gt9], ignore_index=True)
for col in ["Avg # Dep/Month", "Avg $$ Dep/Month", "Average of OD Limit", "Average of Swipes"]:
    nsf_deps_p[col] = nsf_deps_p[col].round(2)
print(f"  {nsf_deps_p.shape[0]} rows")


# ── ANALYSIS 10: Business NSF + Deposits + Swipes ──
print("── Analysis 10: Business NSF + Deposits + Swipes ──")
b10 = df[BUSINESS_OPEN].copy()
b10["NSF Bin"] = pd.cut(b10["Total Items"], bins=NSF_BINS, labels=NSF_LABELS)

nsf_deps_b = b10.groupby("NSF Bin", observed=True).agg(
    **{
        "# of Accounts": ("AcctNo", "count"),
        "Total OD/NSF Items": ("Total Items", "sum"),
        "# of Items Paid": ("Paid Items", "sum"),
        "Avg # Dep/Month": ("Deposit Count", "mean"),
        "Avg $$ Dep/Month": ("Deposit Amount", "mean"),
        "Average of OD Limit": ("OD Limit", "mean"),
        "Average of Swipes": ("Swipes", "mean"),
    }
).reset_index()
nsf_deps_b["NSF Bin"] = nsf_deps_b["NSF Bin"].astype(str)

ta10 = nsf_deps_b["# of Accounts"].sum()
ti10 = nsf_deps_b["Total OD/NSF Items"].sum()
tp10 = nsf_deps_b["# of Items Paid"].sum()
nsf_deps_b["% of Accounts"] = (nsf_deps_b["# of Accounts"] / ta10 * 100).round(2)
nsf_deps_b["% of Items Presented"] = (nsf_deps_b["Total OD/NSF Items"] / ti10 * 100).round(2) if ti10 > 0 else 0
nsf_deps_b["% Pay Rate"] = np.where(nsf_deps_b["Total OD/NSF Items"] > 0, (nsf_deps_b["# of Items Paid"] / nsf_deps_b["Total OD/NSF Items"] * 100).round(1), 0)

gt10 = pd.DataFrame([{
    "NSF Bin": "Grand Total", "# of Accounts": ta10, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti10, "% of Items Presented": 100.0,
    "# of Items Paid": tp10, "% Pay Rate": round(tp10 / ti10 * 100, 1) if ti10 > 0 else 0,
    "Avg # Dep/Month": b10["Deposit Count"].mean(),
    "Avg $$ Dep/Month": b10["Deposit Amount"].mean(),
    "Average of OD Limit": b10["OD Limit"].mean(),
    "Average of Swipes": b10["Swipes"].mean(),
}])
nsf_deps_b = pd.concat([nsf_deps_b, gt10], ignore_index=True)
for col in ["Avg # Dep/Month", "Avg $$ Dep/Month", "Average of OD Limit", "Average of Swipes"]:
    nsf_deps_b[col] = nsf_deps_b[col].round(2)
print(f"  {nsf_deps_b.shape[0]} rows")


# ── ANALYSIS 11: Personal OD Status Stratification ──
print("── Analysis 11: Personal OD Status Stratification ──")
p_od = df[PERSONAL_OPEN].copy()

od_status_personal = p_od.groupby("OD Status").agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum"), "# of Items Paid": ("Paid Items", "sum")}
).reset_index()
od_status_personal["OD Status"] = od_status_personal["OD Status"].astype(str)

ta11 = od_status_personal["# of Accounts"].sum()
ti11 = od_status_personal["Total OD/NSF Items"].sum()
tp11 = od_status_personal["# of Items Paid"].sum()
od_status_personal["% of Accounts"] = (od_status_personal["# of Accounts"] / ta11 * 100).round(2)
od_status_personal["% of Items Presented"] = np.where(ti11 > 0, (od_status_personal["Total OD/NSF Items"] / ti11 * 100).round(2), 0)
od_status_personal["Pay Ratio"] = np.where(od_status_personal["Total OD/NSF Items"] > 0, (od_status_personal["# of Items Paid"] / od_status_personal["Total OD/NSF Items"]).round(4), 0)

gt11 = pd.DataFrame([{
    "OD Status": "Grand Total", "# of Accounts": ta11, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti11, "% of Items Presented": 100.0,
    "# of Items Paid": tp11, "Pay Ratio": round(tp11 / ti11, 4) if ti11 > 0 else 0,
}])
od_status_personal = pd.concat([od_status_personal, gt11], ignore_index=True)
print(f"  {od_status_personal.shape[0]} rows")


# ── ANALYSIS 12: Business OD Status Stratification ──
print("── Analysis 12: Business OD Status Stratification ──")
b_od = df[BUSINESS_OPEN].copy()

od_status_business = b_od.groupby("OD Status").agg(
    **{"# of Accounts": ("AcctNo", "count"), "Total OD/NSF Items": ("Total Items", "sum"), "# of Items Paid": ("Paid Items", "sum")}
).reset_index()
od_status_business["OD Status"] = od_status_business["OD Status"].astype(str)

ta12 = od_status_business["# of Accounts"].sum()
ti12 = od_status_business["Total OD/NSF Items"].sum()
tp12 = od_status_business["# of Items Paid"].sum()
od_status_business["% of Accounts"] = (od_status_business["# of Accounts"] / ta12 * 100).round(2)
od_status_business["% of Items Presented"] = np.where(ti12 > 0, (od_status_business["Total OD/NSF Items"] / ti12 * 100).round(2), 0)
od_status_business["Pay Ratio"] = np.where(od_status_business["Total OD/NSF Items"] > 0, (od_status_business["# of Items Paid"] / od_status_business["Total OD/NSF Items"]).round(4), 0)

gt12 = pd.DataFrame([{
    "OD Status": "Grand Total", "# of Accounts": ta12, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti12, "% of Items Presented": 100.0,
    "# of Items Paid": tp12, "Pay Ratio": round(tp12 / ti12, 4) if ti12 > 0 else 0,
}])
od_status_business = pd.concat([od_status_business, gt12], ignore_index=True)
print(f"  {od_status_business.shape[0]} rows")


# ── ANALYSIS 13: Reg E Summary (Personal) ──
print("── Analysis 13: Reg E Summary ──")
p_rege = df[PERSONAL_OPEN].copy()

reg_e_summary = p_rege.groupby("Reg E Flag").agg(**{"# of Accounts": ("AcctNo", "count")}).reset_index()
ta13 = reg_e_summary["# of Accounts"].sum()
reg_e_summary["% of Accounts"] = (reg_e_summary["# of Accounts"] / ta13 * 100).round(2)

gt13 = pd.DataFrame([{"Reg E Flag": "Grand Total", "# of Accounts": ta13, "% of Accounts": 100.0}])
reg_e_summary = pd.concat([reg_e_summary, gt13], ignore_index=True)
print(f"  Reg E values: {reg_e_summary[reg_e_summary['Reg E Flag'] != 'Grand Total']['Reg E Flag'].tolist()}")


# ── ANALYSIS 14: OD Limit Stratification (Personal) ──
print("── Analysis 14: OD Limit Stratification ──")
p_odl = df[PERSONAL_OPEN].copy()

od_limit_summary = p_odl.groupby("OD Limit").agg(
    **{
        "# of Accounts": ("AcctNo", "count"),
        "Total OD/NSF Items": ("Total Items", "sum"),
        "# of Items Paid": ("Paid Items", "sum"),
        "Avg # Dep/Month": ("Deposit Count", "mean"),
        "Avg $$ Dep/Month": ("Deposit Amount", "mean"),
        "Average of Swipes": ("Swipes", "mean"),
    }
).reset_index()
od_limit_summary["OD Limit"] = od_limit_summary["OD Limit"].astype(int).astype(str)

ta14 = od_limit_summary["# of Accounts"].sum()
ti14 = od_limit_summary["Total OD/NSF Items"].sum()
tp14 = od_limit_summary["# of Items Paid"].sum()
od_limit_summary["% of Accounts"] = (od_limit_summary["# of Accounts"] / ta14 * 100).round(2)
od_limit_summary["% of Items"] = np.where(ti14 > 0, (od_limit_summary["Total OD/NSF Items"] / ti14 * 100).round(2), 0)
od_limit_summary["Pay Ratio"] = np.where(od_limit_summary["Total OD/NSF Items"] > 0, (od_limit_summary["# of Items Paid"] / od_limit_summary["Total OD/NSF Items"]).round(4), 0)

gt14 = pd.DataFrame([{
    "OD Limit": "Grand Total", "# of Accounts": ta14, "% of Accounts": 100.0,
    "Total OD/NSF Items": ti14, "% of Items": 100.0,
    "# of Items Paid": tp14, "Pay Ratio": round(tp14 / ti14, 4) if ti14 > 0 else 0,
    "Avg # Dep/Month": p_odl["Deposit Count"].mean(),
    "Avg $$ Dep/Month": p_odl["Deposit Amount"].mean(),
    "Average of Swipes": p_odl["Swipes"].mean(),
}])
od_limit_summary = pd.concat([od_limit_summary, gt14], ignore_index=True)
od_limit_summary = od_limit_summary[["OD Limit", "# of Accounts", "% of Accounts", "Total OD/NSF Items", "% of Items", "# of Items Paid", "Pay Ratio", "Avg # Dep/Month", "Avg $$ Dep/Month", "Average of Swipes"]]
print(f"  {od_limit_summary.shape[0]} rows")


# ── ANALYSIS 15: Historical Reg E by Year Opened ──
print("── Analysis 15: Historical Reg E by Year Opened ──")
p_hist = df[PERSONAL_OPEN].copy()
p_hist = p_hist[p_hist["Open Date"].notna()].copy()

def assign_year_bin(year):
    if pd.isna(year): return "Unknown"
    y = int(year)
    if y < 2010: return "<2010"
    if y <= 2025: return str(y)
    return "2025+"

p_hist["Year Bin"] = p_hist["Year Opened"].apply(assign_year_bin)

pivot_raw = p_hist.groupby(["Year Bin", "Reg E Flag"], dropna=False).agg(**{"# of Accounts": ("AcctNo", "count")}).reset_index()
pivot_table = pivot_raw.pivot(index="Year Bin", columns="Reg E Flag", values="# of Accounts").fillna(0)

reg_e_flags = sorted([c for c in pivot_table.columns if c is not None])
pivot_table = pivot_table.reindex(columns=reg_e_flags)
pivot_table["# of Accounts"] = pivot_table.sum(axis=1)

# Opt-in %
opt_in_flag = "Y" if "Y" in reg_e_flags else (reg_e_flags[0] if reg_e_flags else None)
if opt_in_flag:
    denom = pivot_table["# of Accounts"].replace({0: pd.NA})
    pivot_table["Opt In %"] = (pivot_table[opt_in_flag] / denom * 100).fillna(0)
else:
    pivot_table["Opt In %"] = 0.0

# Grand Total
grand_total_hist = pd.DataFrame(pivot_table.sum(numeric_only=True)).T
grand_total_hist.index = ["Grand Total"]
if opt_in_flag:
    grand_total_hist["Opt In %"] = (grand_total_hist[opt_in_flag] / grand_total_hist["# of Accounts"] * 100).fillna(0)
pivot_table = pd.concat([pivot_table, grand_total_hist])

# Clean types
int_cols = [c for c in pivot_table.columns if c in reg_e_flags or "# of" in c]
for c in int_cols:
    pivot_table[c] = pivot_table[c].round(0).astype(int)
pivot_table["Opt In %"] = pivot_table["Opt In %"].round(1)

pivot_table = pivot_table.reset_index()
if "Year Bin" in pivot_table.columns:
    pivot_table.rename(columns={"Year Bin": "Year Opened"}, inplace=True)
elif "index" in pivot_table.columns:
    pivot_table.rename(columns={"index": "Year Opened"}, inplace=True)

sort_order = ["<2010"] + [str(y) for y in range(2010, 2026)] + ["2025+", "Unknown", "Grand Total"]
pivot_table["Year Opened"] = pd.Categorical(pivot_table["Year Opened"], categories=sort_order, ordered=True)
pivot_table = pivot_table.sort_values("Year Opened")
print(f"  {pivot_table.shape[0]} rows, Reg E flags: {reg_e_flags}")


# ═══════════════════════════════════════════════════════════════════
# 5. GENERATE POWERPOINT
# ═══════════════════════════════════════════════════════════════════
print("\n═══ Generating PowerPoint ═══")
prs = Presentation()

# Title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "ILS Kickoff Analysis"
title_slide.placeholders[1].text = f"Client 1774 · {datetime.now().strftime('%B %Y')}"

# Slide 1 - Account Status
add_slide_with_table(prs, "Account Status Analysis – All Accounts", stat_code_summary)

# Slide 2 - Account Type
add_slide_with_table(prs, "Account Type Analysis – Open Accounts Only", acct_type_summary)

# Slide 3 - Personal Deposits
add_slide_with_table(prs, "Personal Account Deposit Analysis", personal_deposit_summary)

# Slide 4 - Business Deposits
add_slide_with_table(prs, "Business Account Deposit Analysis", business_deposit_summary)

# Slide 5 - Personal NSF Strat (Volume)
add_slide_with_table(prs, "Personal NSF/OD Stratification – Volume", nsf_strat1)

# Slide 6 - Business NSF Strat (Volume)
add_slide_with_table(prs, "Business NSF/OD Stratification – Volume", nsf_strat_biz)

# Slide 7 - Personal NSF + Pay Ratio
add_slide_with_table(prs, "Personal NSF/OD Stratification – Pay Ratio", nsf_pay_p)

# Slide 8 - Business NSF + Pay Ratio
add_slide_with_table(prs, "Business NSF/OD Stratification – Pay Ratio", nsf_pay_b)

# Slide 9 - Personal NSF + Deposits + Swipes
add_slide_with_table(prs, "Personal NSF/OD – Full Behavioral Metrics", nsf_deps_p)

# Slide 10 - Business NSF + Deposits + Swipes
add_slide_with_table(prs, "Business NSF/OD – Full Behavioral Metrics", nsf_deps_b)

# Slide 11 - Personal OD Status
add_slide_with_table(prs, "Personal OD Status Code Stratification", od_status_personal)

# Slide 12 - Business OD Status
add_slide_with_table(prs, "Business OD Status Code Stratification", od_status_business)

# Slide 13 - Reg E Summary
add_slide_with_table(prs, "Reg E Distribution – Personal Open Accounts", reg_e_summary)

# Slide 14 - OD Limit Stratification
add_slide_with_table(prs, "Personal OD Limit Stratification", od_limit_summary)

# Slide 15 - Historical Reg E
add_slide_with_table(prs, "Historical Reg E Opt-In by Year Opened", pivot_table)

# Save
pptx_path = "1774_ILS_Kickoff_Presentation.pptx"
prs.save(pptx_path)
print(f"✓ PowerPoint saved: {pptx_path}")


# ═══════════════════════════════════════════════════════════════════
# 6. GENERATE EXCEL REPORT
# ═══════════════════════════════════════════════════════════════════
print("\n═══ Generating Excel Report ═══")

now = datetime.now()
month_abbr = calendar.month_abbr[now.month].upper()
year_short = str(now.year)[-2:]
excel_path = Path(f"1774_ILS_Kickoff_Report_{month_abbr}{year_short}.xlsx")

# Create workbook with cover sheet
wb = Workbook()
ws = wb.active
ws.title = "Report Info"
ws["A1"] = "ILS Kickoff Report"
ws["A1"].font = Font(size=20, bold=True)
ws["A3"] = "Report Details:"
ws["A3"].font = Font(size=14, bold=True)

report_info = [
    ("Client ID:", "1774"),
    ("Report Date:", now.strftime("%B %d, %Y")),
    ("Source File:", DATA_PATH.name),
    ("Total Accounts:", f"{len(df):,}"),
]
for i, (label, value) in enumerate(report_info, start=4):
    ws[f"A{i}"] = label
    ws[f"A{i}"].font = Font(bold=True)
    ws[f"B{i}"] = value

ws.column_dimensions["A"].width = 20
ws.column_dimensions["B"].width = 40
wb.save(excel_path)
wb.close()

# Write each analysis as a tab
tabs = [
    ("Stat_Code_Analysis", stat_code_summary),
    ("Account_Type", acct_type_summary),
    ("Personal_Deposits", personal_deposit_summary),
    ("Business_Deposits", business_deposit_summary),
    ("NSF_Strat_Personal", nsf_strat1),
    ("NSF_Strat_Business", nsf_strat_biz),
    ("NSF_PayRatio_Personal", nsf_pay_p),
    ("NSF_PayRatio_Business", nsf_pay_b),
    ("NSF_Full_Personal", nsf_deps_p),
    ("NSF_Full_Business", nsf_deps_b),
    ("OD_Status_Personal", od_status_personal),
    ("OD_Status_Business", od_status_business),
    ("Reg_E_Summary", reg_e_summary),
    ("OD_Limit_Strat", od_limit_summary),
    ("Historical_Reg_E", pivot_table),
]

for sheet_name, data_df in tabs:
    with pd.ExcelWriter(str(excel_path), engine="openpyxl", mode="a", if_sheet_exists="replace") as writer:
        data_df.to_excel(writer, sheet_name=sheet_name, index=False)
    print(f"  ✓ {sheet_name}")

print(f"✓ Excel saved: {excel_path}")

# Cleanup temp chart files
for f in chart_files:
    if os.path.exists(f):
        os.remove(f)

print("\n═══ DONE ═══")
