"""Tests for PowerPoint report generation."""

from pathlib import Path

import pytest
from pptx import Presentation

from ils_kickoff.analyses import run_all_analyses
from ils_kickoff.charts import create_charts
from ils_kickoff.data_loader import load_data
from ils_kickoff.pipeline import PipelineResult
from ils_kickoff.pptx_report import write_pptx_report


@pytest.fixture
def pipeline_result(sample_settings):
    """Run the full pipeline for report testing."""
    df = load_data(sample_settings)
    analyses = run_all_analyses(df, sample_settings)
    charts = create_charts(analyses, sample_settings)
    sample_settings.output_dir.mkdir(parents=True, exist_ok=True)
    return PipelineResult(
        settings=sample_settings, df=df, analyses=analyses, charts=charts,
    )


class TestPptxReport:
    """Test PowerPoint report generation."""

    def test_creates_file(self, pipeline_result, tmp_path):
        path = tmp_path / "test_report.pptx"
        write_pptx_report(pipeline_result, path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_has_title_slide(self, pipeline_result, tmp_path):
        path = tmp_path / "test_report.pptx"
        write_pptx_report(pipeline_result, path)
        prs = Presentation(str(path))
        assert len(prs.slides) >= 1
        # First slide should be title
        first_slide = prs.slides[0]
        assert len(first_slide.shapes) > 0

    def test_has_analysis_slides(self, pipeline_result, tmp_path):
        path = tmp_path / "test_report.pptx"
        write_pptx_report(pipeline_result, path)
        prs = Presentation(str(path))
        successful = [a for a in pipeline_result.analyses if a.error is None]
        # Title slide + at least one slide per analysis
        assert len(prs.slides) >= len(successful) + 1

    def test_slides_have_tables(self, pipeline_result, tmp_path):
        path = tmp_path / "test_report.pptx"
        write_pptx_report(pipeline_result, path)
        prs = Presentation(str(path))
        # Skip title slide (index 0), check analysis slides have tables
        tables_found = 0
        for i, slide in enumerate(prs.slides):
            if i == 0:
                continue
            for shape in slide.shapes:
                if shape.has_table:
                    tables_found += 1
                    break
        assert tables_found > 0
